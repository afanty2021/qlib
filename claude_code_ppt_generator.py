#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Claude Code介绍 PPT生成器
创建一个关于Claude Code的完整演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_claude_code_presentation():
    """创建Claude Code介绍演示文稿"""

    # 创建新的演示文稿
    prs = Presentation()

    # 设置演示文稿的尺寸（16:9宽屏）
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 定义颜色主题（专业技术演示风格）
    primary_color = RGBColor(0, 32, 96)       # 深蓝色 - 主要颜色
    secondary_color = RGBColor(0, 120, 212)   # 亮蓝色 - 次要颜色
    accent_color = RGBColor(255, 127, 0)      # 橙色 - 强调色
    text_color = RGBColor(50, 50, 50)         # 深灰色 - 正文颜色

    # 1. 封面页
    create_title_slide(prs, primary_color, secondary_color, accent_color)

    # 2. 概述页
    create_overview_slide(prs, primary_color, secondary_color, text_color)

    # 3. 核心功能特性页
    create_features_slide(prs, primary_color, secondary_color, accent_color, text_color)

    # 4. 使用场景和应用页
    create_use_cases_slide(prs, primary_color, secondary_color, text_color)

    # 5. 工具和命令介绍页
    create_tools_slide(prs, primary_color, secondary_color, text_color)

    # 6. 最佳实践页
    create_best_practices_slide(prs, primary_color, secondary_color, accent_color, text_color)

    # 7. 总结和Q&A页
    create_summary_slide(prs, primary_color, secondary_color, accent_color)

    # 保存演示文稿
    output_file = "Claude_Code_介绍.pptx"
    prs.save(output_file)
    print(f"演示文稿已成功创建: {output_file}")
    print(f"文件大小: {os.path.getsize(output_file)} 字节")

def create_title_slide(prs, primary_color, secondary_color, accent_color):
    """创建封面页"""
    # 使用标题幻灯片布局
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # 设置标题
    title = slide.shapes.title
    title.text = "Claude Code 介绍"

    # 设置副标题
    subtitle = slide.placeholders[1]
    subtitle.text = "AI驱动的智能编程助手"

    # 格式化标题
    title_fill = title.text_frame.paragraphs[0].font.color
    title_fill.rgb = primary_color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)

    # 格式化副标题
    subtitle_fill = subtitle.text_frame.paragraphs[0].font.color
    subtitle_fill.rgb = secondary_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

def create_overview_slide(prs, primary_color, secondary_color, text_color):
    """创建概述页"""
    # 使用标题和内容布局
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "什么是 Claude Code？"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加内容
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # 清除默认文本

    # 添加定义
    p = tf.add_paragraph()
    p.text = "Claude Code 是 Anthropic 开发的官方命令行界面工具"
    p.font.size = Pt(20)
    p.font.color.rgb = text_color
    p.level = 0

    # 添加核心定位
    p = tf.add_paragraph()
    p.text = "核心定位：为开发者提供AI驱动的编程辅助和代码管理能力"
    p.font.size = Pt(18)
    p.font.color.rgb = secondary_color
    p.level = 0

    # 添加主要特点
    p = tf.add_paragraph()
    p.text = "主要特点："
    p.font.size = Pt(18)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.level = 0

    features = [
        "基于Claude 3.5 Sonnet模型的强大编程能力",
        "本地代码库理解和上下文感知",
        "智能代码生成、重构和优化建议",
        "多项目管理和团队协作支持",
        "与现有开发工作流的深度集成"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.level = 1

def create_features_slide(prs, primary_color, secondary_color, accent_color, text_color):
    """创建核心功能特性页"""
    # 使用标题和内容布局
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "核心功能特性"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加功能特性
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    features = [
        {
            "title": "智能代码分析",
            "description": "深度理解代码库结构，提供智能分析和建议"
        },
        {
            "title": "多文件操作",
            "description": "同时处理多个文件，支持跨文件的代码重构和优化"
        },
        {
            "title": "上下文感知",
            "description": "基于项目结构提供相关的代码建议和最佳实践"
        },
        {
            "title": "实时协作",
            "description": "支持团队协作，共享编程知识和最佳实践"
        },
        {
            "title": "自动化工作流",
            "description": "集成CI/CD流程，自动化代码审查和测试生成"
        },
        {
            "title": "扩展性支持",
            "description": "支持自定义技能和插件，满足特定项目需求"
        }
    ]

    for i, feature in enumerate(features):
        # 添加特性标题
        p = tf.add_paragraph()
        p.text = f"{i+1}. {feature['title']}"
        p.font.size = Pt(18)
        p.font.color.rgb = secondary_color
        p.font.bold = True
        p.level = 0

        # 添加特性描述
        p = tf.add_paragraph()
        p.text = f"   {feature['description']}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.level = 1

def create_use_cases_slide(prs, primary_color, secondary_color, text_color):
    """创建使用场景和应用页"""
    # 使用标题和内容布局
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "使用场景和实际应用"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加使用场景
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    # 开发阶段
    p = tf.add_paragraph()
    p.text = "🔧 开发阶段应用"
    p.font.size = Pt(18)
    p.font.color.rgb = secondary_color
    p.font.bold = True
    p.level = 0

    development_cases = [
        "新功能快速原型设计和实现",
        "代码重构和架构优化建议",
        "技术债务识别和解决方案",
        "单元测试和集成测试生成"
    ]

    for case in development_cases:
        p = tf.add_paragraph()
        p.text = f"• {case}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.level = 1

    # 代码审查阶段
    p = tf.add_paragraph()
    p.text = "🔍 代码审查和维护"
    p.font.size = Pt(18)
    p.font.color.rgb = secondary_color
    p.font.bold = True
    p.level = 0

    review_cases = [
        "自动化代码审查和问题检测",
        "性能瓶颈分析和优化建议",
        "安全漏洞扫描和修复指导",
        "代码质量评估和改进方案"
    ]

    for case in review_cases:
        p = tf.add_paragraph()
        p.text = f"• {case}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.level = 1

def create_tools_slide(prs, primary_color, secondary_color, text_color):
    """创建工具和命令介绍页"""
    # 使用标题和内容布局
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "工具和命令简介"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加命令工具
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    tools = [
        {
            "category": "核心命令",
            "commands": [
                "/help - 显示帮助信息和可用命令",
                "/clear - 清除当前会话历史",
                "/files - 查看项目文件结构",
                "/read <file> - 读取指定文件内容"
            ]
        },
        {
            "category": "代码操作",
            "commands": [
                "/edit <file> - 编辑文件内容",
                "/search <pattern> - 搜索代码模式",
                "/refactor <scope> - 代码重构建议",
                "/analyze <component> - 组件深度分析"
            ]
        },
        {
            "category": "项目管理",
            "commands": [
                "/init <project> - 初始化项目上下文",
                "/status - 显示项目状态和进度",
                "/docs <topic> - 生成项目文档",
                "/test <scope> - 运行或生成测试"
            ]
        }
    ]

    for tool in tools:
        # 添加分类标题
        p = tf.add_paragraph()
        p.text = f"📋 {tool['category']}"
        p.font.size = Pt(18)
        p.font.color.rgb = secondary_color
        p.font.bold = True
        p.level = 0

        # 添加命令列表
        for command in tool['commands']:
            p = tf.add_paragraph()
            p.text = f"• {command}"
            p.font.size = Pt(15)
            p.font.color.rgb = text_color
            p.level = 1

def create_best_practices_slide(prs, primary_color, secondary_color, accent_color, text_color):
    """创建最佳实践页"""
    # 使用标题和内容布局
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "最佳实践建议"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加最佳实践
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    practices = [
        {
            "title": "项目准备",
            "tips": [
                "确保代码库有清晰的结构和文档",
                "提供详细的项目上下文和需求说明",
                "设置合适的代码规范和最佳实践指南"
            ]
        },
        {
            "title": "有效协作",
            "tips": [
                "明确任务范围和预期成果",
                "提供具体的代码示例和参考",
                "及时反馈和调整AI的建议"
            ]
        },
        {
            "title": "质量保证",
            "tips": [
                "定期审查AI生成的代码质量",
                "建立代码审查和测试流程",
                "关注安全性和性能优化建议"
            ]
        },
        {
            "title": "持续优化",
            "tips": [
                "记录成功案例和最佳实践",
                "定制适合项目的工作流程",
                "持续学习和适应新的功能特性"
            ]
        }
    ]

    for practice in practices:
        # 添加实践标题
        p = tf.add_paragraph()
        p.text = f"💡 {practice['title']}"
        p.font.size = Pt(17)
        p.font.color.rgb = secondary_color
        p.font.bold = True
        p.level = 0

        # 添加实践建议
        for tip in practice['tips']:
            p = tf.add_paragraph()
            p.text = f"• {tip}"
            p.font.size = Pt(15)
            p.font.color.rgb = text_color
            p.level = 1

def create_summary_slide(prs, primary_color, secondary_color, accent_color):
    """创建总结和Q&A页"""
    # 使用仅标题布局
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)

    # 设置标题
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "总结与问答"

    # 格式化标题
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # 添加总结文本框
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(11)
    height = Inches(3)

    summary_box = slide.shapes.add_textbox(left, top, width, height)
    tf = summary_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # 添加总结要点
    p = tf.add_paragraph()
    p.text = "📌 核心要点"
    p.font.size = Pt(20)
    p.font.color.rgb = secondary_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    summary_points = [
        "Claude Code 是AI驱动的智能编程助手，显著提升开发效率",
        "提供代码分析、重构、测试生成等全方位编程支持",
        "深度集成现有开发工作流，支持团队协作",
        "通过最佳实践实现代码质量和开发效率的双重提升"
    ]

    for point in summary_points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0

    # 添加Q&A部分
    left = Inches(1)
    top = Inches(4.8)
    width = Inches(11)
    height = Inches(2)

    qa_box = slide.shapes.add_textbox(left, top, width, height)
    tf = qa_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.add_paragraph()
    p.text = "❓ 常见问题"
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    qa_items = [
        "如何在不同项目中配置Claude Code？",
        "Claude Code支持哪些编程语言和框架？",
        "如何确保AI生成代码的安全性？",
        "团队协作的最佳实践是什么？"
    ]

    for item in qa_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0

    # 添加联系信息
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(11)
    height = Inches(0.8)

    contact_box = slide.shapes.add_textbox(left, top, width, height)
    tf = contact_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.add_paragraph()
    p.text = "📧 更多信息：anthropic.com | 📚 文档：docs.anthropic.com"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_claude_code_presentation()